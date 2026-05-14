import os
import json
import yaml
from pptx import Presentation
from pptx.util import Inches

def run_automation():
    # Load Config
    with open("config.yaml", 'r') as f:
        config = yaml.safe_load(f)
    
    # Load Metadata
    with open(config['paths']['metadata_source'], 'r') as f:
        meta = json.load(f)

    # 1. Open Template and Replace Text
    prs = Presentation(config['paths']['template_path'])
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        # Logic to replace placeholders
                        if "{{SN}}" in run.text:
                            run.text = run.text.replace("{{SN}}", meta['sn'])
                        if "{{MODEL}}" in run.text:
                            run.text = run.text.replace("{{MODEL}}", meta['model'])

    # 2. Batch Process Images
    img_folder = config['paths']['image_source']
    images = [f for f in os.listdir(img_folder) if f.endswith(('.jpg', '.png', '.jpeg'))]

    for img_name in images:
        # Layout 6 is usually a blank slide in standard templates
        slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(slide_layout)
        
        img_path = os.path.join(img_folder, img_name)
        
        # Auto-center and Resize
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(config['settings']['image_width_inches'])
        slide.shapes.add_picture(img_path, left, top, width=width)

    # 3. Save Output
    output_path = os.path.join(config['paths']['output_folder'], f"Report_{meta['sn']}.pptx")
    prs.save(output_path)
    print(f"Success: Report generated at {output_path}")

if __name__ == "__main__":
    run_automation()