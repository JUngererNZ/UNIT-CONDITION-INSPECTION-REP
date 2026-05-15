"""
FML Freight Solutions — Inspection Report Generator
====================================================
Usage:
    python generate_report.py                    # process ALL inspections in config.inspections_dir
    python generate_report.py --ba BA3036        # process a single BA number
    python generate_report.py --config my.yaml   # use a custom config file

Output: one .pptx per inspection written to config.output_dir
"""

import argparse
import json
import os
import sys
import glob
from pathlib import Path
from datetime import date

import yaml
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ── Brand colours ──────────────────────────────────────────────────────────────
DARK_TEAL   = RGBColor(0x00, 0x3D, 0x3B)   # slide background
MID_TEAL    = RGBColor(0x00, 0x7A, 0x7A)   # accents
LIGHT_TEAL  = RGBColor(0x5F, 0xC4, 0xBF)   # secondary text / logo
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF0, 0xF4, 0xF4)
DARK_GREY   = RGBColor(0x1A, 0x1A, 0x1A)
MID_GREY    = RGBColor(0x55, 0x55, 0x55)
LIGHT_GREY  = RGBColor(0xCC, 0xCC, 0xCC)
RED_DAMAGE  = RGBColor(0xCC, 0x22, 0x22)
GREEN_OK    = RGBColor(0x22, 0x88, 0x44)
AMBER       = RGBColor(0xE6, 0x8A, 0x00)

# ── Slide dimensions (16:9 widescreen) ─────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ═══════════════════════════════════════════════════════════════════════════════
# Helper utilities
# ═══════════════════════════════════════════════════════════════════════════════

def _set_bg(slide, colour: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour


def _add_textbox(slide, left, top, width, height, text,
                 font_size=12, bold=False, italic=False,
                 colour=WHITE, align=PP_ALIGN.LEFT,
                 wrap=True, font_name="Calibri"):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = colour
    run.font.name = font_name
    return txb


def _add_rect(slide, left, top, width, height, fill_colour, line_colour=None):
    from pptx.util import Pt as UtilPt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_colour
    if line_colour:
        shape.line.color.rgb = line_colour
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def _tick(present: bool, damage: bool) -> str:
    if not present:
        return "✗"
    if damage:
        return "!"          # present but damaged
    return "✓"


def _tick_colour(present: bool, damage: bool) -> RGBColor:
    if not present:
        return RED_DAMAGE
    if damage:
        return AMBER
    return GREEN_OK


def _photos_for_inspection(inspection_dir: str) -> list[str]:
    """Return sorted list of photo paths in inspection_dir/photos/"""
    photo_dir = os.path.join(inspection_dir, "photos")
    if not os.path.isdir(photo_dir):
        return []
    exts = ("*.jpg", "*.jpeg", "*.png", "*.JPG", "*.JPEG", "*.PNG")
    paths = []
    for ext in exts:
        paths.extend(glob.glob(os.path.join(photo_dir, ext)))
    return sorted(set(paths))


# ═══════════════════════════════════════════════════════════════════════════════
# Slide builders
# ═══════════════════════════════════════════════════════════════════════════════

def build_cover_slide(prs: Presentation, data: dict, config: dict, logo_path: str | None):
    layout = prs.slide_layouts[6]   # blank
    slide = prs.slides.add_slide(layout)
    _set_bg(slide, DARK_TEAL)

    # Large company name at bottom
    _add_textbox(slide,
                 Inches(0.3), Inches(5.6), Inches(12.7), Inches(0.9),
                 config["report"]["company_name"],
                 font_size=44, bold=True, colour=LIGHT_TEAL,
                 align=PP_ALIGN.LEFT, font_name="Arial Black")

    _add_textbox(slide,
                 Inches(0.3), Inches(6.45), Inches(12.7), Inches(0.5),
                 config["report"]["tagline"],
                 font_size=13, bold=False, colour=MID_TEAL,
                 align=PP_ALIGN.LEFT, font_name="Calibri Light")

    # Report type label
    _add_textbox(slide,
                 Inches(0.3), Inches(0.2), Inches(12.0), Inches(0.6),
                 "FML WHSE  —  BARTRAC STORAGE UNITS INSPECTION",
                 font_size=16, bold=True, colour=WHITE,
                 align=PP_ALIGN.LEFT, font_name="Calibri")

    # Accent bar
    _add_rect(slide, Inches(0.3), Inches(0.85), Inches(4.0), Inches(0.04), MID_TEAL)

    # Decorative "F" motif (simple lines)
    _add_rect(slide, Inches(5.5),  Inches(1.1), Inches(0.08), Inches(3.2), LIGHT_TEAL)
    _add_rect(slide, Inches(5.5),  Inches(1.1), Inches(1.8),  Inches(0.08), LIGHT_TEAL)
    _add_rect(slide, Inches(5.5),  Inches(2.5), Inches(1.3),  Inches(0.08), MID_TEAL)


def build_checklist_slide(prs: Presentation, data: dict, config: dict):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_bg(slide, OFF_WHITE)

    # ── Header bar ──
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.55), DARK_TEAL)
    _add_textbox(slide,
                 Inches(0.25), Inches(0.08), Inches(12.0), Inches(0.45),
                 "FML WHSE  BARTRAC UNIT CONDITION INSPECTION",
                 font_size=14, bold=True, colour=WHITE,
                 align=PP_ALIGN.LEFT, font_name="Calibri")

    # ── Unit meta block ──
    meta_y = Inches(0.65)
    meta_lines = [
        (f"BA NUMBER:  {data['ba_number']}  —  {data['unit_description']}",   13, True),
        (f"S/N:  {data['serial_number']}",                                     12, False),
        (f"VIN:  {data['vin']}",                                               12, False),
        (f"INSPECTION DATE:  {data['inspection_date']}",                       12, True),
    ]
    if data.get("status_note"):
        meta_lines.append((data["status_note"], 10, False))

    for line_text, fsize, is_bold in meta_lines:
        _add_textbox(slide,
                     Inches(0.25), meta_y, Inches(12.5), Inches(0.28),
                     line_text,
                     font_size=fsize, bold=is_bold, colour=DARK_GREY,
                     align=PP_ALIGN.LEFT)
        meta_y += Inches(0.27)

    # ── Checklist table ──
    table_top   = meta_y + Inches(0.1)
    table_left  = Inches(0.25)
    table_w     = Inches(12.8)
    row_h       = Inches(0.265)

    col_widths  = [Inches(2.6), Inches(1.1), Inches(1.1), Inches(7.8)]
    headers     = ["CHECK POINT", "OFFLOADED /\nSTORAGE", "DAMAGE\nNOTED", "COMMENTS"]

    # Header row
    x = table_left
    for i, (hdr, cw) in enumerate(zip(headers, col_widths)):
        _add_rect(slide, x, table_top, cw, row_h, DARK_TEAL)
        _add_textbox(slide, x + Inches(0.04), table_top + Inches(0.02),
                     cw - Inches(0.08), row_h - Inches(0.04),
                     hdr, font_size=8, bold=True, colour=WHITE,
                     align=PP_ALIGN.CENTER, font_name="Calibri")
        x += cw

    checklist_items = config.get("checklist_items", list(data["checklist"].keys()))

    for row_idx, item_name in enumerate(checklist_items):
        item = data["checklist"].get(item_name, {"present": False, "damage": False, "comments": ""})
        present = item.get("present", False)
        damage  = item.get("damage", False)
        comment = item.get("comments", "")

        row_top = table_top + row_h * (row_idx + 1)
        bg = WHITE if row_idx % 2 == 0 else OFF_WHITE

        x = table_left
        cell_data = [
            (item_name, DARK_GREY, PP_ALIGN.LEFT, 8),
            (_tick(present, damage), _tick_colour(present, damage), PP_ALIGN.CENTER, 11),
            ("!" if damage else ("" if not present else ""), RED_DAMAGE if damage else DARK_GREY, PP_ALIGN.CENTER, 11),
            (comment, MID_GREY, PP_ALIGN.LEFT, 7.5),
        ]
        for (txt, fc, al, fs), cw in zip(cell_data, col_widths):
            _add_rect(slide, x, row_top, cw, row_h, bg, LIGHT_GREY)
            _add_textbox(slide, x + Inches(0.04), row_top + Inches(0.02),
                         cw - Inches(0.08), row_h - Inches(0.04),
                         txt, font_size=fs, colour=fc, align=al, wrap=True)
            x += cw

    # ── General comments box ──
    gc_top = table_top + row_h * (len(checklist_items) + 1) + Inches(0.08)
    _add_rect(slide, table_left, gc_top, table_w, Inches(0.55), RGBColor(0xE8, 0xF4, 0xF4), LIGHT_GREY)
    _add_textbox(slide, table_left + Inches(0.05), gc_top + Inches(0.03),
                 table_w - Inches(0.1), Inches(0.22),
                 "GENERAL COMMENTS:", font_size=7.5, bold=True, colour=DARK_TEAL)
    _add_textbox(slide, table_left + Inches(0.05), gc_top + Inches(0.24),
                 table_w - Inches(0.1), Inches(0.28),
                 data.get("general_comments", ""),
                 font_size=8, colour=DARK_GREY, wrap=True)


def build_photo_slides(prs: Presentation, photos: list[str], photos_per_slide: int, ba_number: str):
    """Build one or more 2-column photo grid slides."""
    cols       = 2
    rows       = photos_per_slide // cols
    pad        = Inches(0.12)
    header_h   = Inches(0.45)
    avail_w    = SLIDE_W - pad * (cols + 1)
    avail_h    = SLIDE_H - header_h - pad * (rows + 1)
    cell_w     = avail_w / cols
    cell_h     = avail_h / rows

    for slide_num, batch_start in enumerate(range(0, len(photos), photos_per_slide)):
        batch = photos[batch_start: batch_start + photos_per_slide]

        layout = prs.slide_layouts[6]
        slide  = prs.slides.add_slide(layout)
        _set_bg(slide, DARK_GREY)

        # Thin header
        _add_rect(slide, Inches(0), Inches(0), SLIDE_W, header_h, DARK_TEAL)
        label = f"{ba_number}  —  CONDITION PHOTOS  ({slide_num + 1})"
        _add_textbox(slide, Inches(0.2), Inches(0.08), Inches(12.0), Inches(0.32),
                     label, font_size=12, bold=True, colour=WHITE,
                     align=PP_ALIGN.LEFT, font_name="Calibri")

        for i, photo_path in enumerate(batch):
            row = i // cols
            col = i %  cols
            left = pad + col * (cell_w + pad)
            top  = header_h + pad + row * (cell_h + pad)

            try:
                pic = slide.shapes.add_picture(photo_path, left, top, cell_w, cell_h)
                # Keep aspect ratio: shrink to fit inside cell
                img_w = pic.width
                img_h = pic.height
                ratio = min(cell_w / img_w, cell_h / img_h)
                pic.width  = int(img_w * ratio)
                pic.height = int(img_h * ratio)
                # Centre in cell
                pic.left = int(left + (cell_w - pic.width)  / 2)
                pic.top  = int(top  + (cell_h - pic.height) / 2)
            except Exception as e:
                print(f"  ⚠  Could not insert photo {photo_path}: {e}")


# ═══════════════════════════════════════════════════════════════════════════════
# Main generation logic
# ═══════════════════════════════════════════════════════════════════════════════

def generate_report(inspection_dir: str, config: dict, output_dir: str):
    data_path = os.path.join(inspection_dir, "inspection_data.json")
    if not os.path.exists(data_path):
        print(f"  ✗  No inspection_data.json found in {inspection_dir} — skipping.")
        return

    with open(data_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    ba = data["ba_number"]
    print(f"  → Generating report for {ba} ...")

    photos           = _photos_for_inspection(inspection_dir)
    photos_per_slide = int(config["report"].get("photos_per_slide", 6))
    logo_path        = config["paths"].get("logo") or None
    if logo_path and not os.path.exists(logo_path):
        logo_path = None

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    build_cover_slide(prs, data, config, logo_path)
    build_checklist_slide(prs, data, config)
    if photos:
        build_photo_slides(prs, photos, photos_per_slide, ba)
    else:
        print(f"     ⚠  No photos found in {inspection_dir}/photos/")

    filename_pattern = config["report"].get("output_filename_pattern",
                                             "{ba_number}_CONDITION_INSPECTION_REPORT.pptx")
    out_filename = filename_pattern.format(ba_number=ba,
                                           date=date.today().strftime("%Y%m%d"))
    out_path = os.path.join(output_dir, out_filename)
    os.makedirs(output_dir, exist_ok=True)
    prs.save(out_path)
    print(f"     ✓  Saved → {out_path}")
    return out_path


def run(config_path: str, ba_filter: str | None = None):
    with open(config_path, "r", encoding="utf-8") as f:
        config = yaml.safe_load(f)

    base_dir        = os.path.dirname(os.path.abspath(config_path))
    inspections_dir = os.path.join(base_dir, config["paths"]["inspections_dir"])
    output_dir      = os.path.join(base_dir, config["paths"]["output_dir"])

    if ba_filter:
        targets = [os.path.join(inspections_dir, ba_filter)]
    else:
        targets = [
            os.path.join(inspections_dir, d)
            for d in sorted(os.listdir(inspections_dir))
            if os.path.isdir(os.path.join(inspections_dir, d))
        ]

    if not targets:
        print("No inspection folders found.")
        sys.exit(1)

    print(f"\nFML Inspection Report Generator")
    print(f"{'─' * 50}")
    results = []
    for t in targets:
        result = generate_report(t, config, output_dir)
        if result:
            results.append(result)

    print(f"\n✓ Done. {len(results)} report(s) generated.\n")
    return results


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="FML Inspection Report Generator")
    parser.add_argument("--config", default="config.yaml",  help="Path to config.yaml")
    parser.add_argument("--ba",     default=None,            help="Generate a single BA number only")
    args = parser.parse_args()

    config_path = args.config
    if not os.path.isabs(config_path):
        config_path = os.path.join(os.path.dirname(__file__), config_path)

    run(config_path, ba_filter=args.ba)
