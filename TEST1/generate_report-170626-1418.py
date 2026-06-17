"""
FML Freight Solutions — Inspection Report Generator
====================================================
Usage:
    python generate_report.py                    # process ALL inspections in config.inspections_dir
    python generate_report.py --ba BA3036        # process a single BA number
    python generate_report.py --config my.yaml   # use a custom config file

Output: one .pptx per inspection written to config.output_dir
"""

import argparse
import io
import json
import os
import sys
import glob
from pathlib import Path
from datetime import date

import yaml
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand colours ──────────────────────────────────────────────────────────────
DARK_TEAL   = RGBColor(0x00, 0x3D, 0x3B)   # slide background (checklist / photo slides)
COVER_BG    = RGBColor(0x00, 0x2C, 0x2B)   # cover background — matched to FML logo pixel (0,44,43)
MID_TEAL    = RGBColor(0x00, 0x7A, 0x7A)   # accents
LIGHT_TEAL  = RGBColor(0x5F, 0xC4, 0xBF)   # secondary text / logo
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF0, 0xF4, 0xF4)
DARK_GREY   = RGBColor(0x1A, 0x1A, 0x1A)
MID_GREY    = RGBColor(0x55, 0x55, 0x55)
LIGHT_GREY  = RGBColor(0xCC, 0xCC, 0xCC)
RED_DAMAGE  = RGBColor(0xCC, 0x22, 0x22)
GREEN_OK    = RGBColor(0x22, 0x88, 0x44)
AMBER       = RGBColor(0xE6, 0x8A, 0x00)

# ── Slide dimensions (portrait A4-ratio) ───────────────────────────────────────
SLIDE_W = Inches(7.5)
SLIDE_H = Inches(13.33)

# ── EMU → pixel conversion at 96 DPI (standard screen resolution) ─────────────
EMU_PER_PX = 9525


# ═══════════════════════════════════════════════════════════════════════════════
# Helper utilities
# ═══════════════════════════════════════════════════════════════════════════════

def _set_bg(slide, colour: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour


def _add_textbox(slide, left, top, width, height, text,
                 font_size=12, bold=False, italic=False,
                 colour=WHITE, align=PP_ALIGN.LEFT,
                 wrap=True, font_name="Calibri"):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = colour
    run.font.name = font_name
    return txb


def _add_rect(slide, left, top, width, height, fill_colour, line_colour=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_colour
    if line_colour:
        shape.line.color.rgb = line_colour
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def _tick(present: bool, damage: bool) -> str:
    if not present: return "✗"
    if damage: return "!"
    return "✓"


def _tick_colour(present: bool, damage: bool) -> RGBColor:
    if not present: return RED_DAMAGE
    if damage: return AMBER
    return GREEN_OK


def _photos_for_inspection(inspection_dir: str, vin: str = None, pin: str = None) -> list[str]:
    """
    Get all photos, with VIN/PIN priority photo as the VERY FIRST image.
    The priority photo will appear in the top-left position of the first photo slide.
    """
    photo_dir = os.path.join(inspection_dir, "photos")
    if not os.path.isdir(photo_dir):
        return []
    
    exts = ("*.jpg", "*.jpeg", "*.png", "*.JPG", "*.JPEG", "*.PNG")
    all_photos = []
    for ext in exts:
        all_photos.extend(glob.glob(os.path.join(photo_dir, ext)))
    all_photos = sorted(set(all_photos))
    
    if not all_photos:
        return []
    
    # Find the VIN/PIN photo - this will be moved to position 0
    priority_photo = None
    photo_dir_lower = photo_dir.lower()
    
    for photo in all_photos:
        photo_name = os.path.basename(photo).upper()
        
        # Check if this photo contains VIN or PIN in the filename
        if vin and vin.upper() in photo_name:
            priority_photo = photo
            print(f"  🎯 Found VIN photo: {os.path.basename(photo)}")
            break
        elif pin and pin.upper() in photo_name:
            priority_photo = photo
            print(f"  🎯 Found PIN photo: {os.path.basename(photo)}")
            break
    
    # If no exact VIN/PIN match, look for keywords
    if not priority_photo:
        keywords = ['VIN', 'PIN', 'SERIAL', 'ID', 'PLATE', 'CHASSIS', 'FRAME', 'IDENT']
        for photo in all_photos:
            photo_name = os.path.basename(photo).upper()
            for keyword in keywords:
                if keyword in photo_name:
                    priority_photo = photo
                    print(f"  🎯 Found ID photo (contains '{keyword}'): {os.path.basename(photo)}")
                    break
            if priority_photo:
                break
    
    # Move priority photo to the VERY FRONT (index 0) - this makes it the first image in the grid
    if priority_photo and priority_photo in all_photos:
        all_photos.remove(priority_photo)
        all_photos.insert(0, priority_photo)
        print(f"  ⭐ Priority photo placed as FIRST image (top-left position)")
    else:
        print(f"  ℹ️  No VIN/PIN photo found - using default ordering")
    
    return all_photos


def _safe_photo_stream(photo_path: str, cell_w_emu: int, cell_h_emu: int,
                       quality: int = 82, scale_factor: float = 1.5) -> io.BytesIO:
    max_w = int((cell_w_emu / EMU_PER_PX) * scale_factor)
    max_h = int((cell_h_emu / EMU_PER_PX) * scale_factor)
    img = Image.open(photo_path)
    if img.format == "MPO": img.seek(0)
    if img.mode in ("RGBA", "LA", "P"): img = img.convert("RGB")
    elif img.mode != "RGB": img = img.convert("RGB")
    try:
        from PIL import ImageOps
        img = ImageOps.exif_transpose(img)
    except Exception: pass
    img.thumbnail((max_w, max_h), Image.LANCZOS)
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=quality, optimize=True)
    buf.seek(0)
    return buf


# ═══════════════════════════════════════════════════════════════════════════════
# Slide builders
# ═══════════════════════════════════════════════════════════════════════════════

def build_cover_slide(prs: Presentation, data: dict, config: dict, logo_path: str | None):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_bg(slide, COVER_BG)

    logo_w    = Inches(6.2)
    logo_left = (SLIDE_W - logo_w) / 2
    logo_top  = Inches(3.8)

    if logo_path and os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, int(logo_left), int(logo_top), int(logo_w))
    else:
        _add_rect(slide, Inches(2.8), Inches(1.5), Inches(0.08), Inches(4.5), LIGHT_TEAL)
        _add_rect(slide, Inches(2.8), Inches(1.5), Inches(2.5),  Inches(0.08), LIGHT_TEAL)
        _add_rect(slide, Inches(2.8), Inches(3.7), Inches(1.8),  Inches(0.08), MID_TEAL)
        _add_textbox(slide, Inches(0.3), Inches(6.3), Inches(6.8), Inches(0.9),
                     config["report"]["company_name"], font_size=28, bold=True,
                     colour=LIGHT_TEAL, align=PP_ALIGN.CENTER, font_name="Arial Black")
        _add_textbox(slide, Inches(0.3), Inches(7.2), Inches(6.8), Inches(0.5),
                     config["report"]["tagline"], font_size=11, bold=False,
                     colour=MID_TEAL, align=PP_ALIGN.CENTER, font_name="Calibri Light")

    # ── Horizontal rule below logo ──
    rule_top = Inches(8.0)
    _add_rect(slide, Inches(0.4), rule_top, SLIDE_W - Inches(0.8), Inches(0.03), MID_TEAL)

    # ── Report details below rule ──
    _add_textbox(slide, Inches(0.4), Inches(8.2), Inches(6.6), Inches(0.5),
                 "BARTRAC UNIT CONDITION INSPECTION REPORT",
                 font_size=11, bold=True, colour=WHITE,
                 align=PP_ALIGN.CENTER, font_name="Calibri")

    _add_textbox(slide, Inches(0.4), Inches(8.85), Inches(6.6), Inches(0.65),
                 f"BA NUMBER:  {data.get('ba_number', '')}",
                 font_size=22, bold=True, colour=LIGHT_TEAL,
                 align=PP_ALIGN.CENTER, font_name="Arial Black")

    _add_textbox(slide, Inches(0.4), Inches(9.6), Inches(6.6), Inches(0.4),
                 data.get("unit_description", ""),
                 font_size=12, bold=False, colour=OFF_WHITE,
                 align=PP_ALIGN.CENTER, font_name="Calibri")

    # ── NEW: Serial Number block added right underneath unit_description ──
    _add_textbox(slide, Inches(0.4), Inches(10.0), Inches(6.6), Inches(0.35),
                 f"SERIAL NUMBER:  {data.get('serial_number', '')}",
                 font_size=11, bold=False, colour=WHITE,
                 align=PP_ALIGN.CENTER, font_name="Calibri")

    # ── Shifted Downward: Inspection Date field moved to accommodate the new line ──
    _add_textbox(slide, Inches(0.4), Inches(10.45), Inches(6.6), Inches(0.35),
                 f"INSPECTION DATE:  {data.get('inspection_date', '')}",
                 font_size=10, bold=False, colour=MID_TEAL,
                 align=PP_ALIGN.CENTER, font_name="Calibri")

    # ── Bottom accent bar ──
    _add_rect(slide, Inches(0), Inches(12.9), SLIDE_W, Inches(0.43), DARK_TEAL)
    _add_textbox(slide, Inches(0.3), Inches(12.93), Inches(6.8), Inches(0.35),
                 config["report"].get("tagline", ""),
                 font_size=8, bold=False, colour=LIGHT_TEAL,
                 align=PP_ALIGN.CENTER, font_name="Calibri Light")


def build_checklist_slide(prs: Presentation, data: dict, config: dict):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_bg(slide, OFF_WHITE)
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.55), DARK_TEAL)
    _add_textbox(slide, Inches(0.25), Inches(0.08), Inches(6.8), Inches(0.45),
                 "FML WHSE  BARTRAC UNIT CONDITION INSPECTION", font_size=12, bold=True, colour=WHITE, align=PP_ALIGN.LEFT, font_name="Calibri")
    meta_y = Inches(0.65)
    meta_lines = [(f"BA NUMBER:  {data['ba_number']}  —  {data['unit_description']}", 12, True), (f"S/N:  {data['serial_number']}", 11, False), (f"VIN:  {data['vin']}", 11, False), (f"INSPECTION DATE:  {data['inspection_date']}", 11, True)]
    if data.get("status_note"): meta_lines.append((data["status_note"], 10, False))
    for line_text, fsize, is_bold in meta_lines:
        _add_textbox(slide, Inches(0.25), meta_y, Inches(6.8), Inches(0.28), line_text, font_size=fsize, bold=is_bold, colour=DARK_GREY, align=PP_ALIGN.LEFT)
        meta_y += Inches(0.27)
    table_top, table_left, table_w, row_h = meta_y + Inches(0.1), Inches(0.25), Inches(7.0), Inches(0.265)
    col_widths, headers = [Inches(1.44), Inches(0.61), Inches(0.61), Inches(4.34)], ["CHECK POINT", "OFFLOADED /\nSTORAGE", "DAMAGE\nNOTED", "COMMENTS"]
    x = table_left
    for i, (hdr, cw) in enumerate(zip(headers, col_widths)):
        _add_rect(slide, x, table_top, cw, row_h, DARK_TEAL)
        _add_textbox(slide, x + Inches(0.04), table_top + Inches(0.02), cw - Inches(0.08), row_h - Inches(0.04), hdr, font_size=7, bold=True, colour=WHITE, align=PP_ALIGN.CENTER, font_name="Calibri")
        x += cw
    checklist_items = config.get("checklist_items", list(data["checklist"].keys()))
    for row_idx, item_name in enumerate(checklist_items):
        item = data["checklist"].get(item_name, {"present": False, "damage": False, "comments": ""})
        row_top, bg = table_top + row_h * (row_idx + 1), WHITE if row_idx % 2 == 0 else OFF_WHITE
        x = table_left
        cell_data = [(item_name, DARK_GREY, PP_ALIGN.LEFT, 7), (_tick(item.get("present", False), item.get("damage", False)), _tick_colour(item.get("present", False), item.get("damage", False)), PP_ALIGN.CENTER, 10), ("!" if item.get("damage", False) else ("" if not item.get("present", False) else ""), RED_DAMAGE if item.get("damage", False) else DARK_GREY, PP_ALIGN.CENTER, 10), (item.get("comments", ""), MID_GREY, PP_ALIGN.LEFT, 7)]
        for (txt, fc, al, fs), cw in zip(cell_data, col_widths):
            _add_rect(slide, x, row_top, cw, row_h, bg, LIGHT_GREY)
            _add_textbox(slide, x + Inches(0.04), row_top + Inches(0.02), cw - Inches(0.08), row_h - Inches(0.04), txt, font_size=fs, colour=fc, align=al, wrap=True)
            x += cw
    gc_top = table_top + row_h * (len(checklist_items) + 1) + Inches(0.08)
    _add_rect(slide, table_left, gc_top, table_w, Inches(0.55), RGBColor(0xE8, 0xF4, 0xF4), LIGHT_GREY)
    _add_textbox(slide, table_left + Inches(0.05), gc_top + Inches(0.03), table_w - Inches(0.1), Inches(0.22), "GENERAL COMMENTS:", font_size=7.5, bold=True, colour=DARK_TEAL)
    _add_textbox(slide, table_left + Inches(0.05), gc_top + Inches(0.24), table_w - Inches(0.1), Inches(0.28), data.get("general_comments", ""), font_size=8, colour=DARK_GREY, wrap=True)

def build_photo_slides(prs: Presentation, photos: list[str], photos_per_slide: int, ba_number: str, photo_quality: int = 82, photo_scale_factor: float = 1.5):
    cols, rows = 2, photos_per_slide // 2
    pad, header_h = Inches(0.12), Inches(0.45)
    avail_w, avail_h = SLIDE_W - pad * (cols + 1), SLIDE_H - header_h - pad * (rows + 1)
    cell_w, cell_h = avail_w / cols, avail_h / rows
    for slide_num, batch_start in enumerate(range(0, len(photos), photos_per_slide)):
        batch = photos[batch_start: batch_start + photos_per_slide]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide, DARK_GREY)
        _add_rect(slide, Inches(0), Inches(0), SLIDE_W, header_h, DARK_TEAL)
        _add_textbox(slide, Inches(0.2), Inches(0.08), Inches(7.0), Inches(0.32), f"{ba_number}  —  CONDITION PHOTOS  ({slide_num + 1})", font_size=12, bold=True, colour=WHITE, align=PP_ALIGN.LEFT, font_name="Calibri")
        for i, photo_path in enumerate(batch):
            row, col = i // cols, i % cols
            left, top = pad + col * (cell_w + pad), header_h + pad + row * (cell_h + pad)
            try:
                stream = _safe_photo_stream(photo_path, cell_w_emu=int(cell_w), cell_h_emu=int(cell_h), quality=photo_quality, scale_factor=photo_scale_factor)
                pic = slide.shapes.add_picture(stream, left, top, int(cell_w), int(cell_h))
                ratio = min(cell_w / pic.width, cell_h / pic.height)
                pic.width, pic.height = int(pic.width * ratio), int(pic.height * ratio)
                pic.left, pic.top = int(left + (cell_w - pic.width) / 2), int(top + (cell_h - pic.height) / 2)
            except Exception as e: print(f"  ⚠  Could not insert photo {photo_path}: {e}")

def build_blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, DARK_GREY)
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.45), DARK_TEAL)
    _add_textbox(slide, Inches(0.2), Inches(0.08), Inches(7.0), Inches(0.32), "NOTES  /  ADDITIONAL INFORMATION", font_size=12, bold=True, colour=WHITE, align=PP_ALIGN.LEFT, font_name="Calibri")
    _add_textbox(slide, Inches(0), Inches(6.4), SLIDE_W, Inches(0.4), "This page intentionally left blank.", font_size=10, bold=False, italic=True, colour=MID_GREY, align=PP_ALIGN.CENTER, font_name="Calibri")
    return slide

def generate_report(inspection_dir: str, config: dict, output_dir: str):
    data_path = os.path.join(inspection_dir, "inspection_data.json")
    if not os.path.exists(data_path):
        print(f"  ⚠  No inspection_data.json found in {inspection_dir}, skipping...")
        return None
    
    with open(data_path, "r", encoding="utf-8") as f:
        data = json.load(f)
    
    ba = data["ba_number"]
    vin = data.get("vin", "")
    pin = data.get("pin", "")  # Get PIN if present
    
    print(f"  📋 VIN from data: {vin}")
    print(f"  📋 PIN from data: {pin}")
    
    # Get photos with VIN/PIN priority photo as FIRST image
    photos = _photos_for_inspection(inspection_dir, vin=vin, pin=pin)
    
    if photos:
        print(f"  📸 Total photos found: {len(photos)}")
        print(f"  📸 FIRST image (top-left position): {os.path.basename(photos[0])}")
    else:
        print(f"  📸 No photos found for {ba}")
    
    photos_per_slide = int(config["report"].get("photos_per_slide", 6))
    photo_quality = int(config["report"].get("photo_quality", 82))
    photo_scale_factor = float(config["report"].get("photo_scale_factor", 1.5))
    logo_path = config["paths"].get("logo") if os.path.exists(config["paths"].get("logo", "")) else None
    
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    
    build_cover_slide(prs, data, config, logo_path)
    build_checklist_slide(prs, data, config)
    build_blank_slide(prs)
    build_blank_slide(prs)
    
    if photos:
        build_photo_slides(prs, photos, photos_per_slide, ba, 
                          photo_quality=photo_quality, 
                          photo_scale_factor=photo_scale_factor)
    
    out_filename = config["report"].get("output_filename_pattern", "{ba_number}_CONDITION_INSPECTION_REPORT.pptx").format(ba_number=ba, date=date.today().strftime("%Y%m%d"))
    out_path = os.path.join(output_dir, out_filename)
    os.makedirs(output_dir, exist_ok=True)
    prs.save(out_path)
    print(f"  ✅ Report saved: {out_path}")
    return out_path

def run(config_path: str, ba_filter: str | None = None):
    print(f"📋 Loading configuration from: {config_path}")
    with open(config_path, "r", encoding="utf-8") as f:
        config = yaml.safe_load(f)
    
    base_dir = os.path.dirname(os.path.abspath(config_path))
    inspections_dir = os.path.join(base_dir, config["paths"]["inspections_dir"])
    output_dir = os.path.join(base_dir, config["paths"]["output_dir"])
    
    print(f"📁 Inspections directory: {inspections_dir}")
    print(f"📁 Output directory: {output_dir}")
    
    if ba_filter:
        targets = [os.path.join(inspections_dir, ba_filter)]
        if not os.path.exists(targets[0]):
            print(f"❌ Error: BA directory '{ba_filter}' not found in {inspections_dir}")
            sys.exit(1)
    else:
        targets = [os.path.join(inspections_dir, d) for d in sorted(os.listdir(inspections_dir)) 
                  if os.path.isdir(os.path.join(inspections_dir, d))]
    
    if not targets:
        print("❌ No inspection directories found.")
        sys.exit(1)
    
    print(f"\n📊 Processing {len(targets)} inspection(s)...\n")
    
    successful = 0
    for t in targets:
        ba_name = os.path.basename(t)
        print(f"🔍 Processing: {ba_name}")
        result = generate_report(t, config, output_dir)
        if result:
            successful += 1
        print()
    
    print(f"✅ Completed! {successful}/{len(targets)} reports generated successfully.")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="FML Freight Solutions - Inspection Report Generator")
    parser.add_argument("--config", default="config.yaml", help="Path to config file (default: config.yaml)")
    parser.add_argument("--ba", default=None, help="Process a single BA number (folder name)")
    args = parser.parse_args()
    
    config_path = args.config if os.path.isabs(args.config) else os.path.join(os.path.dirname(__file__), args.config)
    
    if not os.path.exists(config_path):
        print(f"❌ Error: Config file not found: {config_path}")
        sys.exit(1)
    
    run(config_path, ba_filter=args.ba)